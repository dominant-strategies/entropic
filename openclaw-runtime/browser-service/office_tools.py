#!/usr/bin/env python3

import csv
import io
import json
import os
import posixpath
import re
import shutil
import sys
import tempfile
import zipfile
from typing import Dict, Iterable, List, Optional, Tuple
from xml.etree import ElementTree as ET

try:
    from openpyxl import Workbook, load_workbook
except ImportError:  # pragma: no cover - runtime dependency may be absent in dev
    Workbook = None
    load_workbook = None

try:
    from pptx import Presentation as PptxPresentation
    from pptx.util import Inches
except ImportError:  # pragma: no cover - runtime dependency may be absent in dev
    PptxPresentation = None
    Inches = None

WORKSPACE_ROOT = os.environ.get("ENTROPIC_WORKSPACE_PATH", "/data/workspace")

AIO_SPEC = "agent-interpretable-object"
AIO_ROOT = "Agent-Intepretable-Object"
AIO_MANIFEST = f"{AIO_ROOT}/manifest.yaml"
AIO_CORE_KEYWORDS = f"{AIO_ROOT}/core/keywords.yaml"
AIO_CORE_ONTOLOGY = f"{AIO_ROOT}/core/ontology.yaml"
AIO_CORE_CANONICAL_INSTANCE = f"{AIO_ROOT}/core/canonical-instance.yaml"
AIO_FAMILY_TABULAR = f"{AIO_ROOT}/families/tabular-space.yaml"
AIO_FAMILY_LINEAR_TEXT = f"{AIO_ROOT}/families/linear-text.yaml"
AIO_FAMILY_SLIDE_SPACE = f"{AIO_ROOT}/families/slide-space.yaml"
AIO_FAMILY_GRAPHICS_SCENE = f"{AIO_ROOT}/families/graphics-scene.md"
AIO_KIND_SPREADSHEET = f"{AIO_ROOT}/kinds/spreadsheet.yaml"
AIO_KIND_DOCUMENT = f"{AIO_ROOT}/kinds/document.yaml"
AIO_KIND_PRESENTATION = f"{AIO_ROOT}/kinds/presentation.yaml"

XLSX_NS = "http://schemas.openxmlformats.org/spreadsheetml/2006/main"
DOCX_NS = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
PPTX_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
DRAWING_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
DOC_REL_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
PKG_REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
CONTENT_TYPES_NS = "http://schemas.openxmlformats.org/package/2006/content-types"

ET.register_namespace("", XLSX_NS)
ET.register_namespace("a", DRAWING_NS)
ET.register_namespace("p", PPTX_NS)
ET.register_namespace("r", DOC_REL_NS)
ET.register_namespace("w", DOCX_NS)


def xlsx_tag(name: str) -> str:
    return f"{{{XLSX_NS}}}{name}"


def docx_tag(name: str) -> str:
    return f"{{{DOCX_NS}}}{name}"


def pptx_tag(name: str) -> str:
    return f"{{{PPTX_NS}}}{name}"


def drawing_tag(name: str) -> str:
    return f"{{{DRAWING_NS}}}{name}"


def local_name(tag: str) -> str:
    return tag.rsplit("}", 1)[-1] if "}" in tag else tag


def normalize_workspace_root() -> str:
    return posixpath.normpath(WORKSPACE_ROOT)


def resolve_workspace_path(raw_path: str) -> str:
    if not isinstance(raw_path, str) or not raw_path.strip():
        raise ValueError("A workspace file path is required.")
    root = normalize_workspace_root()
    trimmed = raw_path.strip()
    if trimmed == root or trimmed.startswith(f"{root}/"):
        full_path = posixpath.normpath(trimmed)
    else:
        normalized_relative = posixpath.normpath(f"/{trimmed}").lstrip("/")
        full_path = posixpath.join(root, normalized_relative)
    if full_path != root and not full_path.startswith(f"{root}/"):
        raise ValueError("The requested path is outside /data/workspace.")
    return full_path


def path_metadata(path: str) -> Dict[str, object]:
    exists = os.path.exists(path)
    if not exists:
        return {
            "exists": False,
            "modifiedMs": None,
            "size": 0,
            "etag": "missing",
        }
    stats = os.stat(path)
    modified_ms = int(round(stats.st_mtime_ns / 1_000_000))
    return {
        "exists": True,
        "modifiedMs": modified_ms,
        "size": int(stats.st_size),
        "etag": f"{modified_ms}:{stats.st_size}",
    }


def assert_expected_etag(path: str, expected_etag: Optional[str]) -> None:
    if not expected_etag:
        return
    current = path_metadata(path)["etag"]
    if current != expected_etag:
        raise RuntimeError(
            "The file changed on disk while you were editing it. Reload the viewer before saving again."
        )


def ensure_parent_dir(path: str) -> None:
    os.makedirs(posixpath.dirname(path), exist_ok=True)


def atomic_write_bytes(path: str, data: bytes) -> None:
    ensure_parent_dir(path)
    temp_fd, temp_path = tempfile.mkstemp(prefix=".entropic-office-", dir=posixpath.dirname(path))
    try:
        with os.fdopen(temp_fd, "wb") as handle:
            handle.write(data)
        os.replace(temp_path, path)
    finally:
        try:
            if os.path.exists(temp_path):
                os.unlink(temp_path)
        except OSError:
            pass


def atomic_write_text(path: str, text: str) -> None:
    atomic_write_bytes(path, text.encode("utf-8"))


def split_extension(path: str) -> str:
    return posixpath.splitext(path)[1].lower()


def basename_without_extension(path: str) -> str:
    return posixpath.splitext(posixpath.basename(path))[0] or "object"


def slugify_identifier(value: object, fallback: str) -> str:
    text = str(value or "").strip().lower()
    text = re.sub(r"[^a-z0-9]+", "-", text).strip("-")
    return text or fallback


def aio_definition_refs(kind: str) -> Dict[str, object]:
    families: List[str] = []
    kind_path = ""
    if kind == "spreadsheet":
        families = [AIO_FAMILY_TABULAR]
        kind_path = AIO_KIND_SPREADSHEET
    elif kind == "document":
        families = [AIO_FAMILY_LINEAR_TEXT]
        kind_path = AIO_KIND_DOCUMENT
    elif kind == "presentation":
        families = [AIO_FAMILY_SLIDE_SPACE, AIO_FAMILY_GRAPHICS_SCENE]
        kind_path = AIO_KIND_PRESENTATION
    return {
        "manifest": AIO_MANIFEST,
        "core": [AIO_CORE_KEYWORDS, AIO_CORE_ONTOLOGY, AIO_CORE_CANONICAL_INSTANCE],
        "families": families,
        "kind": kind_path,
    }


def aio_source_block(document: Dict[str, object], extra_notes: Optional[List[str]] = None) -> Dict[str, object]:
    notes = []
    warning = str(document.get("warning") or "").strip()
    if warning:
        notes.append(warning)
    for note in extra_notes or []:
        normalized = str(note or "").strip()
        if normalized:
            notes.append(normalized)
    return {
        "path": str(document.get("path") or ""),
        "exists": bool(document.get("exists")),
        "etag": str(document.get("etag") or "missing"),
        "modifiedMs": document.get("modifiedMs"),
        "size": int(document.get("size") or 0),
        "notes": notes,
    }


def aio_envelope(
    kind: str,
    document: Dict[str, object],
    object_payload: Dict[str, object],
    capability_set: str,
    extra_notes: Optional[List[str]] = None,
) -> Dict[str, object]:
    return {
        "spec": AIO_SPEC,
        "kind": kind,
        "format": str(document.get("format") or "").strip(),
        "definitions": aio_definition_refs(kind),
        "realization": {
            "name": "entropic-office",
            "mode": "automation",
            "capability_set": capability_set,
        },
        "source": aio_source_block(document, extra_notes=extra_notes),
        "object": object_payload,
    }


def extract_aio_source_path(payload: Dict[str, object]) -> str:
    source = payload.get("source")
    if isinstance(source, dict):
        return str(source.get("path") or "").strip()
    return ""


def extract_aio_source_etag(payload: Dict[str, object]) -> Optional[str]:
    source = payload.get("source")
    if not isinstance(source, dict):
        return None
    etag = str(source.get("etag") or "").strip()
    return etag or None


def ensure_aio_payload(payload: Dict[str, object]) -> Dict[str, object]:
    if not isinstance(payload, dict) or payload.get("spec") != AIO_SPEC:
        raise RuntimeError("Expected an agent-interpretable-object payload.")
    return payload


def col_index_to_label(index: int) -> str:
    if index < 1:
        raise ValueError("Column index must be positive.")
    label = []
    value = index
    while value > 0:
        value, remainder = divmod(value - 1, 26)
        label.append(chr(65 + remainder))
    return "".join(reversed(label))


def label_to_col_index(label: str) -> int:
    value = 0
    for char in label.upper():
        if not ("A" <= char <= "Z"):
            raise ValueError(f"Invalid column label `{label}`.")
        value = value * 26 + (ord(char) - 64)
    return value


def cell_ref(row: int, col: int) -> str:
    return f"{col_index_to_label(col)}{row}"


def parse_cell_ref(ref: str) -> Tuple[int, int]:
    match = re.fullmatch(r"([A-Za-z]+)(\d+)", ref or "")
    if not match:
        raise ValueError(f"Invalid cell reference `{ref}`.")
    return int(match.group(2)), label_to_col_index(match.group(1))


def is_numeric_value(value: str) -> bool:
    if not isinstance(value, str):
        return False
    trimmed = value.strip()
    if not trimmed:
        return False
    return re.fullmatch(r"[+-]?(?:\d+(?:\.\d+)?|\.\d+)", trimmed) is not None


def coerce_boolean_text(value: str) -> Optional[str]:
    if not isinstance(value, str):
        return None
    lowered = value.strip().lower()
    if lowered in {"true", "1", "yes"}:
        return "1"
    if lowered in {"false", "0", "no"}:
        return "0"
    return None


def safe_inline_text(parent: ET.Element, value: str) -> None:
    inline = ET.SubElement(parent, xlsx_tag("is"))
    text = ET.SubElement(inline, xlsx_tag("t"))
    text.text = value


def read_shared_strings(archive: zipfile.ZipFile) -> List[str]:
    if "xl/sharedStrings.xml" not in archive.namelist():
        return []
    root = ET.fromstring(archive.read("xl/sharedStrings.xml"))
    values: List[str] = []
    for item in root:
        if local_name(item.tag) != "si":
            continue
        text_parts = []
        for text_node in item.iter():
            if local_name(text_node.tag) == "t":
                text_parts.append(text_node.text or "")
        values.append("".join(text_parts))
    return values


def decode_xlsx_cell(cell: ET.Element, shared_strings: List[str]) -> Dict[str, object]:
    ref = cell.attrib.get("r", "")
    row, col = parse_cell_ref(ref)
    formula_node = cell.find(xlsx_tag("f"))
    value_node = cell.find(xlsx_tag("v"))
    cell_type = cell.attrib.get("t")
    formula = formula_node.text or "" if formula_node is not None and formula_node.text else None
    raw_value = value_node.text or "" if value_node is not None and value_node.text else ""
    if cell_type == "s":
        try:
            display_value = shared_strings[int(raw_value)]
        except (IndexError, ValueError):
            display_value = ""
        kind = "string"
    elif cell_type == "inlineStr":
        text_parts = []
        inline = cell.find(xlsx_tag("is"))
        if inline is not None:
            for node in inline.iter():
                if local_name(node.tag) == "t":
                    text_parts.append(node.text or "")
        display_value = "".join(text_parts)
        kind = "string"
    elif cell_type == "b":
        display_value = "TRUE" if raw_value == "1" else "FALSE"
        kind = "boolean"
    else:
        display_value = raw_value
        kind = "number" if raw_value and is_numeric_value(raw_value) else "string"
    return {
        "ref": ref,
        "row": row,
        "col": col,
        "formula": formula,
        "value": display_value,
        "display": display_value if formula is None else (raw_value or display_value),
        "kind": "formula" if formula else kind,
    }


def parse_xlsx_sheet(sheet_name: str, payload: bytes, shared_strings: List[str]) -> Dict[str, object]:
    root = ET.fromstring(payload)
    sheet_data = root.find(xlsx_tag("sheetData"))
    cells = []
    max_row = 0
    max_col = 0
    if sheet_data is not None:
        for row in sheet_data.findall(xlsx_tag("row")):
            for cell in row.findall(xlsx_tag("c")):
                decoded = decode_xlsx_cell(cell, shared_strings)
                cells.append(decoded)
                max_row = max(max_row, int(decoded["row"]))
                max_col = max(max_col, int(decoded["col"]))
    return {
        "name": sheet_name,
        "cells": cells,
        "rowCount": max_row,
        "colCount": max_col,
    }


def normalize_relationship_target(base_dir: str, target: str) -> str:
    if target.startswith("/"):
        return target.lstrip("/")
    return posixpath.normpath(posixpath.join(base_dir, target))


def read_xlsx_document(path: str) -> Dict[str, object]:
    metadata = path_metadata(path)
    if not metadata["exists"]:
        return {
            "kind": "spreadsheet",
            "format": "xlsx",
            "path": path,
            **metadata,
            "warning": None,
            "sheets": [{"name": "Sheet1", "cells": [], "rowCount": 0, "colCount": 0}],
        }
    with zipfile.ZipFile(path, "r") as archive:
        workbook = ET.fromstring(archive.read("xl/workbook.xml"))
        rels = ET.fromstring(archive.read("xl/_rels/workbook.xml.rels"))
        targets_by_id: Dict[str, str] = {}
        for rel in rels:
            if local_name(rel.tag) != "Relationship":
                continue
            rel_id = rel.attrib.get("Id")
            target = rel.attrib.get("Target")
            if rel_id and target:
                targets_by_id[rel_id] = normalize_relationship_target("xl", target)
        shared_strings = read_shared_strings(archive)
        sheets: List[Dict[str, object]] = []
        sheets_parent = workbook.find(xlsx_tag("sheets"))
        if sheets_parent is not None:
            for sheet in sheets_parent.findall(xlsx_tag("sheet")):
                rel_id = sheet.attrib.get(f"{{{DOC_REL_NS}}}id", "")
                target = targets_by_id.get(rel_id)
                if not target:
                    continue
                if target not in archive.namelist():
                    continue
                sheets.append(
                    parse_xlsx_sheet(
                        sheet.attrib.get("name", f"Sheet{len(sheets) + 1}"),
                        archive.read(target),
                        shared_strings,
                    )
                )
        if not sheets:
            sheets = [{"name": "Sheet1", "cells": [], "rowCount": 0, "colCount": 0}]
        return {
            "kind": "spreadsheet",
            "format": "xlsx",
            "path": path,
            **metadata,
            "warning": None,
            "sheets": sheets,
        }


def read_csv_document(path: str) -> Dict[str, object]:
    metadata = path_metadata(path)
    if not metadata["exists"]:
        return {
            "kind": "spreadsheet",
            "format": "csv",
            "path": path,
            **metadata,
            "warning": None,
            "sheets": [{"name": "Sheet1", "cells": [], "rowCount": 0, "colCount": 0}],
        }
    with open(path, "r", encoding="utf-8", newline="") as handle:
        reader = csv.reader(handle)
        rows = list(reader)
    cells = []
    max_col = 0
    for row_index, row in enumerate(rows, start=1):
        max_col = max(max_col, len(row))
        for col_index, value in enumerate(row, start=1):
            if value == "":
                continue
            cells.append(
                {
                    "ref": cell_ref(row_index, col_index),
                    "row": row_index,
                    "col": col_index,
                    "formula": value[1:] if value.startswith("=") else None,
                    "value": value,
                    "display": value,
                    "kind": "formula" if value.startswith("=") else "string",
                }
            )
    return {
        "kind": "spreadsheet",
        "format": "csv",
        "path": path,
        **metadata,
        "warning": None,
        "sheets": [{"name": "Sheet1", "cells": cells, "rowCount": len(rows), "colCount": max_col}],
    }


def legacy_minimal_xlsx(path: str) -> bool:
    if not os.path.exists(path):
        return False
    try:
        with zipfile.ZipFile(path, "r") as archive:
            entries = set(archive.namelist())
    except (OSError, zipfile.BadZipFile):
        return False
    required_minimal = {
        "[Content_Types].xml",
        "_rels/.rels",
        "xl/workbook.xml",
        "xl/_rels/workbook.xml.rels",
    }
    if not required_minimal.issubset(entries):
        return False
    has_sheet = any(name.startswith("xl/worksheets/") and name.endswith(".xml") for name in entries)
    if not has_sheet:
        return False
    # Older Entropic-generated workbooks omitted these standard parts, which ONLYOFFICE
    # rejects even though the file is readable by the lightweight inspector.
    missing_standard_parts = {
        "docProps/app.xml",
        "docProps/core.xml",
        "xl/styles.xml",
        "xl/theme/theme1.xml",
    }
    return bool(entries.isdisjoint(missing_standard_parts))


def normalize_sheet_payload(sheet: Dict[str, object]) -> Dict[str, object]:
    name = str(sheet.get("name") or "Sheet1").strip() or "Sheet1"
    cleaned_cells = []
    max_row = 0
    max_col = 0
    for raw_cell in sheet.get("cells") or []:
        if not isinstance(raw_cell, dict):
            continue
        raw_ref = str(raw_cell.get("ref") or "").strip().upper()
        if raw_ref:
            row, col = parse_cell_ref(raw_ref)
            ref = raw_ref
        else:
            row = int(raw_cell.get("row") or 0)
            col = int(raw_cell.get("col") or 0)
            if row <= 0 or col <= 0:
                continue
            ref = cell_ref(row, col)
        formula = str(raw_cell.get("formula") or "").strip()
        if formula.startswith("="):
            formula = formula[1:]
        value = str(raw_cell.get("value") or "")
        if formula:
            cleaned_cells.append(
                {
                    "ref": ref,
                    "row": row,
                    "col": col,
                    "formula": formula,
                    "value": value,
                    "display": str(raw_cell.get("display") or value),
                    "kind": "formula",
                }
            )
        elif value != "":
            kind = str(raw_cell.get("kind") or "")
            if kind not in {"number", "boolean", "string"}:
                if is_numeric_value(value):
                    kind = "number"
                elif coerce_boolean_text(value) is not None:
                    kind = "boolean"
                else:
                    kind = "string"
            cleaned_cells.append(
                {
                    "ref": ref,
                    "row": row,
                    "col": col,
                    "formula": None,
                    "value": value,
                    "display": str(raw_cell.get("display") or value),
                    "kind": kind,
                }
            )
        max_row = max(max_row, row)
        max_col = max(max_col, col)
    cleaned_cells.sort(key=lambda cell: (int(cell["row"]), int(cell["col"])))
    return {"name": name, "cells": cleaned_cells, "rowCount": max_row, "colCount": max_col}


def openpyxl_available() -> bool:
    return Workbook is not None and load_workbook is not None


def coerce_openpyxl_value(cell: Dict[str, object]):
    formula = str(cell.get("formula") or "").strip()
    if formula:
        return f"={formula}"
    value = str(cell.get("value") or "")
    kind = str(cell.get("kind") or "string")
    if kind == "boolean":
        return coerce_boolean_text(value) == "1"
    if kind == "number" and is_numeric_value(value):
        trimmed = value.strip()
        if re.fullmatch(r"[+-]?\d+", trimmed):
            try:
                return int(trimmed)
            except ValueError:
                return trimmed
        try:
            return float(trimmed)
        except ValueError:
            return trimmed
    return value


def build_sheet_root(sheet: Dict[str, object], existing_bytes: Optional[bytes]) -> bytes:
    existing_rows: Dict[int, ET.Element] = {}
    existing_cells: Dict[str, ET.Element] = {}
    if existing_bytes:
        root = ET.fromstring(existing_bytes)
    else:
        root = ET.Element(xlsx_tag("worksheet"))
        ET.SubElement(root, xlsx_tag("sheetData"))
    sheet_data = root.find(xlsx_tag("sheetData"))
    if sheet_data is None:
        sheet_data = ET.SubElement(root, xlsx_tag("sheetData"))
    for row in sheet_data.findall(xlsx_tag("row")):
        row_index = int(row.attrib.get("r") or 0)
        if row_index > 0:
            existing_rows[row_index] = row
        for cell in row.findall(xlsx_tag("c")):
            cell_key = cell.attrib.get("r")
            if cell_key:
                existing_cells[cell_key] = cell
    for child in list(sheet_data):
        sheet_data.remove(child)

    rows_by_index: Dict[int, List[Dict[str, object]]] = {}
    for cell in sheet["cells"]:
        rows_by_index.setdefault(int(cell["row"]), []).append(cell)

    for row_index in sorted(rows_by_index):
        row_element = ET.Element(xlsx_tag("row"))
        if row_index in existing_rows:
            row_element.attrib.update(existing_rows[row_index].attrib)
        row_element.attrib["r"] = str(row_index)
        for cell in rows_by_index[row_index]:
            ref = str(cell["ref"])
            existing_cell = existing_cells.get(ref)
            cell_element = ET.Element(xlsx_tag("c"))
            if existing_cell is not None:
                for key, value in existing_cell.attrib.items():
                    if key not in {"r", "t"}:
                        cell_element.attrib[key] = value
            cell_element.attrib["r"] = ref
            formula = cell.get("formula")
            value = str(cell.get("value") or "")
            if formula:
                ET.SubElement(cell_element, xlsx_tag("f")).text = str(formula)
                cached = str(cell.get("display") or value)
                if cached != "":
                    boolean_value = coerce_boolean_text(cached)
                    if boolean_value is not None:
                        cell_element.attrib["t"] = "b"
                        ET.SubElement(cell_element, xlsx_tag("v")).text = boolean_value
                    elif is_numeric_value(cached):
                        ET.SubElement(cell_element, xlsx_tag("v")).text = cached.strip()
                    else:
                        cell_element.attrib["t"] = "str"
                        ET.SubElement(cell_element, xlsx_tag("v")).text = cached
            else:
                kind = str(cell.get("kind") or "string")
                if kind == "boolean":
                    boolean_value = coerce_boolean_text(value) or "0"
                    cell_element.attrib["t"] = "b"
                    ET.SubElement(cell_element, xlsx_tag("v")).text = boolean_value
                elif kind == "number" and is_numeric_value(value):
                    ET.SubElement(cell_element, xlsx_tag("v")).text = value.strip()
                else:
                    cell_element.attrib["t"] = "inlineStr"
                    safe_inline_text(cell_element, value)
            row_element.append(cell_element)
        sheet_data.append(row_element)

    dimension_ref = "A1"
    if sheet["cells"]:
        max_row = max(int(cell["row"]) for cell in sheet["cells"])
        max_col = max(int(cell["col"]) for cell in sheet["cells"])
        dimension_ref = f"A1:{cell_ref(max_row, max_col)}"
    dimension = root.find(xlsx_tag("dimension"))
    if dimension is None:
        dimension = ET.Element(xlsx_tag("dimension"))
        root.insert(0, dimension)
    dimension.attrib["ref"] = dimension_ref
    return ET.tostring(root, encoding="utf-8", xml_declaration=True)


def minimal_xlsx_entries(sheets: List[Dict[str, object]]) -> Dict[str, bytes]:
    workbook = ET.Element(xlsx_tag("workbook"))
    sheets_root = ET.SubElement(workbook, xlsx_tag("sheets"))
    for index, sheet in enumerate(sheets, start=1):
        sheet_el = ET.SubElement(sheets_root, xlsx_tag("sheet"))
        sheet_el.set("name", str(sheet["name"]))
        sheet_el.set("sheetId", str(index))
        sheet_el.set(f"{{{DOC_REL_NS}}}id", f"rId{index}")
    calc_pr = ET.SubElement(workbook, xlsx_tag("calcPr"))
    calc_pr.set("fullCalcOnLoad", "1")
    calc_pr.set("calcMode", "auto")

    rels = ET.Element(f"{{{PKG_REL_NS}}}Relationships")
    for index in range(1, len(sheets) + 1):
        relationship = ET.SubElement(rels, f"{{{PKG_REL_NS}}}Relationship")
        relationship.set("Id", f"rId{index}")
        relationship.set(
            "Type",
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/worksheet",
        )
        relationship.set("Target", f"worksheets/sheet{index}.xml")

    content_types = ET.Element(f"{{{CONTENT_TYPES_NS}}}Types")
    default_rels = ET.SubElement(content_types, f"{{{CONTENT_TYPES_NS}}}Default")
    default_rels.set("Extension", "rels")
    default_rels.set("ContentType", "application/vnd.openxmlformats-package.relationships+xml")
    default_xml = ET.SubElement(content_types, f"{{{CONTENT_TYPES_NS}}}Default")
    default_xml.set("Extension", "xml")
    default_xml.set("ContentType", "application/xml")
    workbook_override = ET.SubElement(content_types, f"{{{CONTENT_TYPES_NS}}}Override")
    workbook_override.set("PartName", "/xl/workbook.xml")
    workbook_override.set(
        "ContentType",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet.main+xml",
    )
    for index in range(1, len(sheets) + 1):
        sheet_override = ET.SubElement(content_types, f"{{{CONTENT_TYPES_NS}}}Override")
        sheet_override.set("PartName", f"/xl/worksheets/sheet{index}.xml")
        sheet_override.set(
            "ContentType",
            "application/vnd.openxmlformats-officedocument.spreadsheetml.worksheet+xml",
        )

    root_rels = ET.Element(f"{{{PKG_REL_NS}}}Relationships")
    workbook_rel = ET.SubElement(root_rels, f"{{{PKG_REL_NS}}}Relationship")
    workbook_rel.set("Id", "rId1")
    workbook_rel.set(
        "Type",
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument",
    )
    workbook_rel.set("Target", "xl/workbook.xml")

    entries = {
        "[Content_Types].xml": ET.tostring(content_types, encoding="utf-8", xml_declaration=True),
        "_rels/.rels": ET.tostring(root_rels, encoding="utf-8", xml_declaration=True),
        "xl/workbook.xml": ET.tostring(workbook, encoding="utf-8", xml_declaration=True),
        "xl/_rels/workbook.xml.rels": ET.tostring(rels, encoding="utf-8", xml_declaration=True),
    }
    for index, sheet in enumerate(sheets, start=1):
        entries[f"xl/worksheets/sheet{index}.xml"] = build_sheet_root(sheet, None)
    return entries


def write_xlsx_document(path: str, sheets: List[Dict[str, object]], expected_etag: Optional[str]) -> Dict[str, object]:
    assert_expected_etag(path, expected_etag)
    normalized_sheets = [normalize_sheet_payload(sheet) for sheet in sheets] or [
        {"name": "Sheet1", "cells": [], "rowCount": 0, "colCount": 0}
    ]

    if openpyxl_available():
        if os.path.exists(path):
            workbook = load_workbook(path)
            if len(workbook.worksheets) != len(normalized_sheets):
                raise RuntimeError(
                    "Adding or removing worksheets is not supported yet. Edit the existing sheets or create a new workbook."
                )
        else:
            workbook = Workbook()
            while len(workbook.worksheets) < len(normalized_sheets):
                workbook.create_sheet()
            while len(workbook.worksheets) > len(normalized_sheets):
                workbook.remove(workbook.worksheets[-1])

        for index, sheet in enumerate(normalized_sheets):
            worksheet = workbook.worksheets[index]
            worksheet.title = str(sheet["name"])
            max_row = max(worksheet.max_row or 0, int(sheet.get("rowCount") or 0))
            max_col = max(worksheet.max_column or 0, int(sheet.get("colCount") or 0))
            for row in range(1, max_row + 1):
                for col in range(1, max_col + 1):
                    worksheet.cell(row=row, column=col).value = None
            for cell in sheet["cells"]:
                worksheet[str(cell["ref"])].value = coerce_openpyxl_value(cell)

        temp_fd, temp_path = tempfile.mkstemp(prefix=".entropic-office-", dir=posixpath.dirname(path))
        os.close(temp_fd)
        try:
            workbook.save(temp_path)
            os.replace(temp_path, path)
        finally:
            try:
                if os.path.exists(temp_path):
                    os.unlink(temp_path)
            except OSError:
                pass
        return read_xlsx_document(path)

    if os.path.exists(path):
        with zipfile.ZipFile(path, "r") as archive:
            entries = {name: archive.read(name) for name in archive.namelist()}
        workbook_root = ET.fromstring(entries["xl/workbook.xml"])
        rels_root = ET.fromstring(entries["xl/_rels/workbook.xml.rels"])
        workbook_sheets = workbook_root.find(xlsx_tag("sheets"))
        if workbook_sheets is None:
            raise RuntimeError("The workbook is missing sheet metadata.")
        existing_names = workbook_sheets.findall(xlsx_tag("sheet"))
        if len(existing_names) != len(normalized_sheets):
            raise RuntimeError(
                "Adding or removing worksheets is not supported yet. Edit the existing sheets or create a new workbook."
            )
        rel_targets: Dict[str, str] = {}
        for rel in rels_root:
            if local_name(rel.tag) != "Relationship":
                continue
            rel_id = rel.attrib.get("Id")
            target = rel.attrib.get("Target")
            if rel_id and target:
                rel_targets[rel_id] = normalize_relationship_target("xl", target)
        for index, sheet in enumerate(existing_names):
            rel_id = sheet.attrib.get(f"{{{DOC_REL_NS}}}id", "")
            target = rel_targets.get(rel_id)
            if not target:
                continue
            sheet.set("name", str(normalized_sheets[index]["name"]))
            entries[target] = build_sheet_root(normalized_sheets[index], entries.get(target))
        entries["xl/workbook.xml"] = ET.tostring(workbook_root, encoding="utf-8", xml_declaration=True)
    else:
        entries = minimal_xlsx_entries(normalized_sheets)

    temp_fd, temp_path = tempfile.mkstemp(prefix=".entropic-office-", dir=posixpath.dirname(path))
    os.close(temp_fd)
    try:
        with zipfile.ZipFile(temp_path, "w", compression=zipfile.ZIP_DEFLATED) as archive:
            for name in sorted(entries):
                archive.writestr(name, entries[name])
        os.replace(temp_path, path)
    finally:
        try:
            if os.path.exists(temp_path):
                os.unlink(temp_path)
        except OSError:
            pass
    return read_xlsx_document(path)


def write_csv_document(path: str, sheets: List[Dict[str, object]], expected_etag: Optional[str]) -> Dict[str, object]:
    assert_expected_etag(path, expected_etag)
    normalized_sheet = normalize_sheet_payload(sheets[0] if sheets else {"name": "Sheet1", "cells": []})
    max_row = max((int(cell["row"]) for cell in normalized_sheet["cells"]), default=0)
    max_col = max((int(cell["col"]) for cell in normalized_sheet["cells"]), default=0)
    grid = [["" for _ in range(max_col)] for _ in range(max_row)]
    for cell in normalized_sheet["cells"]:
        row = int(cell["row"]) - 1
        col = int(cell["col"]) - 1
        if row < 0 or col < 0:
            continue
        if cell.get("formula"):
            grid[row][col] = f"={cell['formula']}"
        else:
            grid[row][col] = str(cell.get("value") or "")
    buffer = io.StringIO()
    writer = csv.writer(buffer)
    for row in grid:
        writer.writerow(row)
    atomic_write_text(path, buffer.getvalue())
    return read_csv_document(path)


def read_docx_document(path: str) -> Dict[str, object]:
    metadata = path_metadata(path)
    if not metadata["exists"]:
        return {
            "kind": "document",
            "format": "docx",
            "path": path,
            **metadata,
            "warning": "Rich formatting and tables are preserved when possible, but the editor currently focuses on paragraph text.",
            "paragraphs": [],
        }
    with zipfile.ZipFile(path, "r") as archive:
        root = ET.fromstring(archive.read("word/document.xml"))
    body = root.find(docx_tag("body"))
    paragraphs: List[str] = []
    if body is not None:
        for child in body:
            if local_name(child.tag) != "p":
                continue
            text_parts = []
            for node in child.iter():
                if local_name(node.tag) == "t":
                    text_parts.append(node.text or "")
            paragraphs.append("".join(text_parts))
    return {
        "kind": "document",
        "format": "docx",
        "path": path,
        **metadata,
        "warning": "Rich formatting and tables are preserved when possible, but the editor currently focuses on paragraph text.",
        "paragraphs": paragraphs,
    }


def build_doc_paragraph(text: str) -> ET.Element:
    paragraph = ET.Element(docx_tag("p"))
    run = ET.SubElement(paragraph, docx_tag("r"))
    text_node = ET.SubElement(run, docx_tag("t"))
    text_node.text = text
    return paragraph


def minimal_docx_entries(paragraphs: List[str]) -> Dict[str, bytes]:
    document = ET.Element(docx_tag("document"))
    body = ET.SubElement(document, docx_tag("body"))
    for paragraph in paragraphs:
        body.append(build_doc_paragraph(paragraph))
    ET.SubElement(body, docx_tag("sectPr"))

    content_types = ET.Element(f"{{{CONTENT_TYPES_NS}}}Types")
    default_rels = ET.SubElement(content_types, f"{{{CONTENT_TYPES_NS}}}Default")
    default_rels.set("Extension", "rels")
    default_rels.set("ContentType", "application/vnd.openxmlformats-package.relationships+xml")
    default_xml = ET.SubElement(content_types, f"{{{CONTENT_TYPES_NS}}}Default")
    default_xml.set("Extension", "xml")
    default_xml.set("ContentType", "application/xml")
    document_override = ET.SubElement(content_types, f"{{{CONTENT_TYPES_NS}}}Override")
    document_override.set("PartName", "/word/document.xml")
    document_override.set(
        "ContentType",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml",
    )

    rels = ET.Element(f"{{{PKG_REL_NS}}}Relationships")
    rel = ET.SubElement(rels, f"{{{PKG_REL_NS}}}Relationship")
    rel.set("Id", "rId1")
    rel.set(
        "Type",
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument",
    )
    rel.set("Target", "word/document.xml")

    return {
        "[Content_Types].xml": ET.tostring(content_types, encoding="utf-8", xml_declaration=True),
        "_rels/.rels": ET.tostring(rels, encoding="utf-8", xml_declaration=True),
        "word/document.xml": ET.tostring(document, encoding="utf-8", xml_declaration=True),
    }


def write_docx_document(path: str, paragraphs: List[str], expected_etag: Optional[str]) -> Dict[str, object]:
    assert_expected_etag(path, expected_etag)
    normalized_paragraphs = [str(value) for value in paragraphs]
    if os.path.exists(path):
        with zipfile.ZipFile(path, "r") as archive:
            entries = {name: archive.read(name) for name in archive.namelist()}
        root = ET.fromstring(entries["word/document.xml"])
        body = root.find(docx_tag("body"))
        if body is None:
            body = ET.SubElement(root, docx_tag("body"))
        section = None
        if len(body) > 0 and local_name(body[-1].tag) == "sectPr":
            section = body[-1]
        template_children = [child for child in list(body) if child is not section]
        body[:] = []
        remaining = list(normalized_paragraphs)
        for child in template_children:
            if local_name(child.tag) == "p":
                if remaining:
                    body.append(build_doc_paragraph(remaining.pop(0)))
            else:
                body.append(child)
        for paragraph in remaining:
            body.append(build_doc_paragraph(paragraph))
        if section is None:
            section = ET.Element(docx_tag("sectPr"))
        body.append(section)
        entries["word/document.xml"] = ET.tostring(root, encoding="utf-8", xml_declaration=True)
    else:
        entries = minimal_docx_entries(normalized_paragraphs)

    temp_fd, temp_path = tempfile.mkstemp(prefix=".entropic-office-", dir=posixpath.dirname(path))
    os.close(temp_fd)
    try:
        with zipfile.ZipFile(temp_path, "w", compression=zipfile.ZIP_DEFLATED) as archive:
            for name in sorted(entries):
                archive.writestr(name, entries[name])
        os.replace(temp_path, path)
    finally:
        try:
            if os.path.exists(temp_path):
                os.unlink(temp_path)
        except OSError:
            pass
    return read_docx_document(path)


PRESENTATION_AUTOMATION_WARNING = (
    "Structured presentation automation edits slide text, notes, layout references, object ordering, table cell text, "
    "image placement, and basic object frames in place. Themes, chart data, and animation semantics are preserved when "
    "possible but are not yet fully modeled."
)

PPTX_TITLE_PLACEHOLDER_TYPES = {"TITLE", "CENTER_TITLE"}
PPTX_BODY_PLACEHOLDER_TYPES = {"BODY", "SUBTITLE", "OBJECT"}
PPTX_NOTES_IGNORED_PLACEHOLDER_TYPES = {
    "HEADER",
    "DATE",
    "FOOTER",
    "SLIDE_NUMBER",
    "HDR",
    "DT",
    "FTR",
    "SLDNUM",
}


def python_pptx_available() -> bool:
    return PptxPresentation is not None and Inches is not None


def emu_frame_dict(x: object, y: object, w: object, h: object) -> Dict[str, object]:
    return {
        "x": int(x or 0),
        "y": int(y or 0),
        "w": int(w or 0),
        "h": int(h or 0),
        "unit": "emu",
    }


def normalize_presentation_frame(frame: object) -> Optional[Dict[str, object]]:
    if not isinstance(frame, dict):
        return None
    if not any(key in frame for key in {"x", "y", "w", "h", "left", "top", "width", "height"}):
        return None
    return emu_frame_dict(
        frame.get("x", frame.get("left", 0)),
        frame.get("y", frame.get("top", 0)),
        frame.get("w", frame.get("width", 0)),
        frame.get("h", frame.get("height", 0)),
    )


def normalize_layout_ref(raw: object) -> Optional[str]:
    text = str(raw or "").strip()
    if not text:
        return None
    if text.startswith("layout:"):
        return text
    return f"layout:{slugify_identifier(text, 'layout')}"


def normalize_master_ref(raw: object) -> Optional[str]:
    text = str(raw or "").strip()
    if not text:
        return None
    if text.startswith("master:"):
        return text
    return f"master:{slugify_identifier(text, 'master')}"


def body_items_from_outline(items: object) -> List[Dict[str, object]]:
    normalized: List[Dict[str, object]] = []
    if not isinstance(items, list):
        return normalized
    for item in items:
        if isinstance(item, dict):
            text = str(item.get("text") or "").strip()
            if not text:
                continue
            level = int(item.get("level") or 0)
        else:
            text = str(item or "").strip()
            if not text:
                continue
            level = 0
        normalized.append({"text": text, "level": max(level, 0)})
    return normalized


def normalize_presentation_body(body: object) -> Dict[str, object]:
    if isinstance(body, dict):
        kind = str(body.get("kind") or "").strip()
        if kind == "outline" or isinstance(body.get("items"), list):
            items = body_items_from_outline(body.get("items"))
            return {"kind": "outline", "items": items}
        bullets = body.get("bullets")
        if isinstance(bullets, list):
            return {
                "kind": "bullets",
                "bullets": [str(value) for value in bullets if str(value or "").strip()],
            }
        text = str(body.get("text") or "").strip()
        return {"kind": "text", "text": text}
    if isinstance(body, (list, tuple)):
        bullets = [str(value) for value in body if str(value or "").strip()]
        return {"kind": "bullets", "bullets": bullets}
    text = str(body or "").strip()
    return {"kind": "text", "text": text}


def presentation_paragraphs_from_body(body: object) -> List[Dict[str, object]]:
    normalized = normalize_presentation_body(body)
    kind = str(normalized.get("kind") or "").strip()
    if kind == "outline":
        return body_items_from_outline(normalized.get("items"))
    if kind == "bullets":
        return [{"text": value, "level": 0} for value in normalized.get("bullets") or []]
    text = str(normalized.get("text") or "").strip()
    return [{"text": line, "level": 0} for line in text.splitlines() if line.strip()]


def presentation_body_from_paragraphs(paragraphs: List[Dict[str, object]]) -> Dict[str, object]:
    cleaned = body_items_from_outline(paragraphs)
    if not cleaned:
        return {"kind": "text", "text": ""}
    if len(cleaned) == 1 and int(cleaned[0].get("level") or 0) == 0:
        return {"kind": "text", "text": str(cleaned[0].get("text") or "")}
    if any(int(item.get("level") or 0) != 0 for item in cleaned):
        return {"kind": "outline", "items": cleaned}
    return {"kind": "bullets", "bullets": [str(item.get("text") or "") for item in cleaned]}


def normalize_presentation_notes(notes: object) -> List[str]:
    if isinstance(notes, list):
        return [str(value) for value in notes if str(value or "").strip()]
    text = str(notes or "").strip()
    return [text] if text else []


def presentation_text_value(body: object) -> str:
    paragraphs = presentation_paragraphs_from_body(body)
    return "\n".join(str(item.get("text") or "") for item in paragraphs if str(item.get("text") or "").strip())


def presentation_body_lines(body: Dict[str, object]) -> List[str]:
    return [str(item.get("text") or "") for item in presentation_paragraphs_from_body(body)]


def normalize_string_matrix(raw: object) -> List[List[str]]:
    rows: List[List[str]] = []
    if not isinstance(raw, list):
        return rows
    for raw_row in raw:
        if not isinstance(raw_row, list):
            continue
        rows.append([str(value or "") for value in raw_row])
    return rows


def normalize_chart_series(raw: object) -> List[Dict[str, object]]:
    normalized: List[Dict[str, object]] = []
    if not isinstance(raw, list):
        return normalized
    for index, item in enumerate(raw, start=1):
        if not isinstance(item, dict):
            continue
        entry: Dict[str, object] = {
            "id": str(item.get("id") or f"series:{index}"),
        }
        name = str(item.get("name") or "").strip()
        if name:
            entry["name"] = name
        values = item.get("values")
        if isinstance(values, list):
            entry["values"] = [str(value) if isinstance(value, bool) else value for value in values]
        if len(entry) > 1:
            normalized.append(entry)
    return normalized


def normalize_presentation_transition(raw: object) -> Optional[Dict[str, object]]:
    if not isinstance(raw, dict):
        return None
    result: Dict[str, object] = {}
    effect = str(raw.get("effect") or "").strip()
    if effect:
        result["effect"] = effect
    speed = str(raw.get("speed") or "").strip()
    if speed:
        result["speed"] = speed
    if raw.get("advance_on_click") is not None:
        result["advance_on_click"] = bool(raw.get("advance_on_click"))
    if raw.get("advance_after_ms") is not None:
        try:
            result["advance_after_ms"] = int(raw.get("advance_after_ms") or 0)
        except (TypeError, ValueError):
            pass
    if raw.get("duration_ms") is not None:
        try:
            result["duration_ms"] = int(raw.get("duration_ms") or 0)
        except (TypeError, ValueError):
            pass
    return result or None


def normalize_presentation_timeline(raw: object) -> Optional[Dict[str, object]]:
    if not isinstance(raw, dict):
        return None
    result: Dict[str, object] = {}
    if raw.get("has_main_sequence") is not None:
        result["has_main_sequence"] = bool(raw.get("has_main_sequence"))
    if raw.get("effect_count") is not None:
        try:
            result["effect_count"] = int(raw.get("effect_count") or 0)
        except (TypeError, ValueError):
            pass
    return result or None


def normalize_presentation_object_payload(
    raw_object: object,
    slide_index: int,
    object_index: int,
) -> Optional[Dict[str, object]]:
    if not isinstance(raw_object, dict):
        return None
    kind = str(raw_object.get("kind") or "").strip() or "shape"
    result: Dict[str, object] = {
        "id": str(raw_object.get("id") or f"object:{slide_index}:{object_index}"),
        "kind": kind,
    }
    name = str(raw_object.get("name") or "").strip()
    if name:
        result["name"] = name
    placeholder_kind = str(raw_object.get("placeholder_kind") or "").strip().upper()
    if placeholder_kind:
        result["placeholder_kind"] = placeholder_kind
    frame = normalize_presentation_frame(raw_object.get("frame"))
    if frame:
        result["frame"] = frame
    if "text" in raw_object:
        result["text"] = str(raw_object.get("text") or "")
    if "body" in raw_object or "bullets" in raw_object or "items" in raw_object:
        result["body"] = normalize_presentation_body(raw_object.get("body", raw_object))
    image_ref = str(raw_object.get("image_ref") or "").strip()
    if image_ref:
        result["image_ref"] = image_ref
    image_name = str(raw_object.get("image_name") or "").strip()
    if image_name:
        result["image_name"] = image_name
    alt_text = str(raw_object.get("alt_text") or "").strip()
    if alt_text:
        result["alt_text"] = alt_text
    chart_kind = str(raw_object.get("chart_kind") or "").strip()
    if chart_kind:
        result["chart_kind"] = chart_kind
    chart_title = str(raw_object.get("title") or "").strip() if kind == "chart" else ""
    if chart_title:
        result["title"] = chart_title
    categories = raw_object.get("categories")
    if isinstance(categories, list):
        result["categories"] = [str(value or "") for value in categories]
    series = normalize_chart_series(raw_object.get("series"))
    if series:
        result["series"] = series
    table = raw_object.get("table")
    if isinstance(table, dict):
        next_table: Dict[str, object] = {}
        if table.get("rows") is not None:
            next_table["rows"] = int(table.get("rows") or 0)
        if table.get("cols") is not None:
            next_table["cols"] = int(table.get("cols") or 0)
        cells = normalize_string_matrix(table.get("cells"))
        if cells:
            next_table["cells"] = cells
            next_table["rows"] = max(int(next_table.get("rows") or 0), len(cells))
            next_table["cols"] = max(
                int(next_table.get("cols") or 0),
                max((len(row) for row in cells), default=0),
            )
        if next_table:
            result["table"] = next_table
    return result


def derive_presentation_title(objects: List[Dict[str, object]]) -> Optional[str]:
    for obj in objects:
        if str(obj.get("kind") or "") == "title" or str(obj.get("placeholder_kind") or "") in PPTX_TITLE_PLACEHOLDER_TYPES:
            if "text" in obj:
                return str(obj.get("text") or "")
            if "body" in obj:
                return presentation_text_value(obj.get("body"))
    return None


def derive_presentation_body(objects: List[Dict[str, object]]) -> Optional[Dict[str, object]]:
    for obj in objects:
        if str(obj.get("placeholder_kind") or "") in PPTX_BODY_PLACEHOLDER_TYPES:
            if "body" in obj:
                return normalize_presentation_body(obj.get("body"))
            if "text" in obj:
                return {"kind": "text", "text": str(obj.get("text") or "")}
    for obj in objects:
        if str(obj.get("kind") or "") == "text_box":
            if "body" in obj:
                return normalize_presentation_body(obj.get("body"))
            if "text" in obj:
                return {"kind": "text", "text": str(obj.get("text") or "")}
    return None


def normalize_presentation_slide_payload(
    raw_slide: object,
    index: int,
    derive_summary: bool = False,
) -> Dict[str, object]:
    result: Dict[str, object] = {
        "id": f"slide:{index}",
        "kind": "slide",
        "index": index,
    }
    if not isinstance(raw_slide, dict):
        return result
    result["id"] = str(raw_slide.get("id") or f"slide:{index}")
    layout_ref = normalize_layout_ref(raw_slide.get("layout_ref") or raw_slide.get("layout"))
    if layout_ref:
        result["layout_ref"] = layout_ref
    master_ref = normalize_master_ref(raw_slide.get("master_ref"))
    if master_ref:
        result["master_ref"] = master_ref
    if "title" in raw_slide:
        result["title"] = str(raw_slide.get("title") or "")
    if "body" in raw_slide:
        result["body"] = normalize_presentation_body(raw_slide.get("body"))
    if "notes" in raw_slide:
        result["notes"] = normalize_presentation_notes(raw_slide.get("notes"))
    transition = normalize_presentation_transition(raw_slide.get("transition"))
    if transition:
        result["transition"] = transition
    timeline = normalize_presentation_timeline(raw_slide.get("timeline"))
    if timeline:
        result["timeline"] = timeline
    remove_object_ids = [
        str(value).strip()
        for value in (raw_slide.get("remove_object_ids") or [])
        if str(value or "").strip()
    ]
    if remove_object_ids:
        result["remove_object_ids"] = remove_object_ids
    if raw_slide.get("reorder_objects") is True:
        result["reorder_objects"] = True
    if raw_slide.get("prune_missing_objects") is True:
        result["prune_missing_objects"] = True
    objects = []
    for object_index, raw_object in enumerate(raw_slide.get("objects") or [], start=1):
        normalized = normalize_presentation_object_payload(raw_object, index, object_index)
        if normalized is not None:
            objects.append(normalized)
    if objects:
        result["objects"] = objects
    if derive_summary and objects:
        if "title" not in result:
            derived_title = derive_presentation_title(objects)
            if derived_title is not None:
                result["title"] = derived_title
        if "body" not in result:
            derived_body = derive_presentation_body(objects)
            if derived_body is not None:
                result["body"] = derived_body
    return result


def drawing_paragraph_texts(node: ET.Element) -> List[str]:
    paragraphs: List[str] = []
    for paragraph in node.findall(f".//{drawing_tag('p')}"):
        text_parts = []
        for text_node in paragraph.iter():
            if local_name(text_node.tag) == "t":
                text_parts.append(text_node.text or "")
        text = "".join(text_parts).strip()
        if text:
            paragraphs.append(text)
    return paragraphs


def xml_transition_payload(slide_root: ET.Element) -> Optional[Dict[str, object]]:
    transition = slide_root.find(pptx_tag("transition"))
    if transition is None:
        return None
    payload: Dict[str, object] = {}
    speed = str(transition.attrib.get("spd") or "").strip()
    if speed:
        payload["speed"] = speed
    adv_click = transition.attrib.get("advClick")
    if adv_click is not None:
        payload["advance_on_click"] = adv_click not in {"0", "false", "False"}
    adv_time = transition.attrib.get("advTm")
    if adv_time:
        try:
            payload["advance_after_ms"] = int(adv_time)
        except ValueError:
            pass
    for child in list(transition):
        name = local_name(child.tag)
        if name != "extLst":
            payload["effect"] = name
            break
    return payload or None


def xml_timeline_payload(slide_root: ET.Element) -> Optional[Dict[str, object]]:
    timing = slide_root.find(pptx_tag("timing"))
    if timing is None:
        return None
    ctn_count = sum(1 for node in timing.iter() if local_name(node.tag) == "cTn")
    return {
        "has_main_sequence": ctn_count > 0,
        "effect_count": max(ctn_count - 1, 0),
    }


def xml_placeholder_type(shape: ET.Element) -> Optional[str]:
    placeholder = shape.find(
        f"{pptx_tag('nvSpPr')}/{pptx_tag('nvPr')}/{pptx_tag('ph')}"
    )
    if placeholder is None:
        return None
    raw_type = str(placeholder.attrib.get("type") or "").strip()
    return (raw_type or "body").upper()


def xml_shape_frame(shape: ET.Element) -> Optional[Dict[str, object]]:
    xfrm = shape.find(f".//{drawing_tag('xfrm')}")
    if xfrm is None:
        return None
    off = xfrm.find(drawing_tag("off"))
    ext = xfrm.find(drawing_tag("ext"))
    if off is None or ext is None:
        return None
    return emu_frame_dict(
        off.attrib.get("x", 0),
        off.attrib.get("y", 0),
        ext.attrib.get("cx", 0),
        ext.attrib.get("cy", 0),
    )


def xml_shape_non_visual(shape: ET.Element) -> Tuple[str, str]:
    paths = (
        (pptx_tag("nvSpPr"), pptx_tag("cNvPr")),
        (pptx_tag("nvPicPr"), pptx_tag("cNvPr")),
        (pptx_tag("nvGraphicFramePr"), pptx_tag("cNvPr")),
    )
    for first, second in paths:
        node = shape.find(f"{first}/{second}")
        if node is not None:
            return (
                str(node.attrib.get("id") or "").strip(),
                str(node.attrib.get("name") or "").strip(),
            )
    return ("", "")


def xml_shape_object(shape: ET.Element, fallback_id: str) -> Optional[Dict[str, object]]:
    object_id_value, name = xml_shape_non_visual(shape)
    object_id = f"shape:{object_id_value or fallback_id}"
    frame = xml_shape_frame(shape)
    placeholder_kind = xml_placeholder_type(shape)

    if local_name(shape.tag) == "sp":
        body = presentation_body_from_paragraphs(
            [{"text": paragraph, "level": 0} for paragraph in drawing_paragraph_texts(shape)]
        )
        if presentation_text_value(body) or name:
            kind = "title" if placeholder_kind in PPTX_TITLE_PLACEHOLDER_TYPES else "text_box"
            result: Dict[str, object] = {"id": object_id, "kind": kind}
            if name:
                result["name"] = name
            if placeholder_kind:
                result["placeholder_kind"] = placeholder_kind
            if frame:
                result["frame"] = frame
            if kind == "title":
                result["text"] = presentation_text_value(body)
            else:
                result["body"] = body
            return result
        return None

    if local_name(shape.tag) == "pic":
        result = {"id": object_id, "kind": "image"}
        if name:
            result["name"] = name
            result["image_name"] = name
        if frame:
            result["frame"] = frame
        return result

    if local_name(shape.tag) == "graphicFrame":
        result: Dict[str, object] = {"id": object_id, "kind": "shape"}
        if name:
            result["name"] = name
        if frame:
            result["frame"] = frame
        if shape.find(f".//{drawing_tag('tbl')}") is not None:
            rows = len(shape.findall(f".//{drawing_tag('tr')}"))
            cols = 0
            cells: List[List[str]] = []
            first_row = shape.find(f".//{drawing_tag('tr')}")
            if first_row is not None:
                cols = len(first_row.findall(drawing_tag("tc")))
            for row in shape.findall(f".//{drawing_tag('tr')}"):
                cell_row: List[str] = []
                for cell in row.findall(drawing_tag("tc")):
                    cell_row.append("\n".join(drawing_paragraph_texts(cell)))
                if cell_row:
                    cells.append(cell_row)
            result["kind"] = "table"
            result["table"] = {"rows": rows, "cols": cols}
            if cells:
                result["table"]["cells"] = cells
            return result
        for node in shape.iter():
            if local_name(node.tag) == "chart":
                result["kind"] = "chart"
                return result
    return None


def read_pptx_notes_xml(archive: zipfile.ZipFile, slide_path: str) -> List[str]:
    rels_path = posixpath.join(
        posixpath.dirname(slide_path),
        "_rels",
        f"{posixpath.basename(slide_path)}.rels",
    )
    if rels_path not in archive.namelist():
        return []
    rels_root = ET.fromstring(archive.read(rels_path))
    notes_target = None
    for rel in rels_root:
        if local_name(rel.tag) != "Relationship":
            continue
        rel_type = str(rel.attrib.get("Type") or "")
        if rel_type.endswith("/notesSlide"):
            target = rel.attrib.get("Target")
            if target:
                notes_target = normalize_relationship_target(posixpath.dirname(slide_path), target)
                break
    if not notes_target or notes_target not in archive.namelist():
        return []
    notes_root = ET.fromstring(archive.read(notes_target))
    notes: List[str] = []
    for shape in notes_root.findall(f".//{pptx_tag('sp')}"):
        placeholder_type = xml_placeholder_type(shape)
        if placeholder_type in PPTX_NOTES_IGNORED_PLACEHOLDER_TYPES:
            continue
        notes.extend(drawing_paragraph_texts(shape))
    return notes


def read_pptx_document_xml(path: str, metadata: Dict[str, object]) -> Dict[str, object]:
    with zipfile.ZipFile(path, "r") as archive:
        presentation_xml = "ppt/presentation.xml"
        presentation_rels_xml = "ppt/_rels/presentation.xml.rels"
        if presentation_xml not in archive.namelist() or presentation_rels_xml not in archive.namelist():
            raise RuntimeError("The presentation is missing required PPTX metadata.")
        presentation_root = ET.fromstring(archive.read(presentation_xml))
        presentation_rels_root = ET.fromstring(archive.read(presentation_rels_xml))
        slide_targets: Dict[str, str] = {}
        for rel in presentation_rels_root:
            if local_name(rel.tag) != "Relationship":
                continue
            rel_id = rel.attrib.get("Id")
            target = rel.attrib.get("Target")
            if rel_id and target:
                slide_targets[rel_id] = normalize_relationship_target("ppt", target)

        slides: List[Dict[str, object]] = []
        slide_id_list = presentation_root.find(pptx_tag("sldIdLst"))
        if slide_id_list is not None:
            for index, slide_ref in enumerate(slide_id_list.findall(pptx_tag("sldId")), start=1):
                slide_stable_id = str(slide_ref.attrib.get("id") or index)
                rel_id = slide_ref.attrib.get(f"{{{DOC_REL_NS}}}id")
                slide_path = slide_targets.get(rel_id or "")
                if not slide_path or slide_path not in archive.namelist():
                    continue
                slide_root = ET.fromstring(archive.read(slide_path))
                objects: List[Dict[str, object]] = []
                for shape_index, shape in enumerate(
                    slide_root.findall(f".//{pptx_tag('sp')}") + slide_root.findall(f".//{pptx_tag('pic')}") + slide_root.findall(f".//{pptx_tag('graphicFrame')}"),
                    start=1,
                ):
                    obj = xml_shape_object(shape, f"{slide_stable_id}:{shape_index}")
                    if obj is None:
                        continue
                    objects.append(obj)
                slides.append(
                    normalize_presentation_slide_payload(
                        {
                            "id": f"slide:{slide_stable_id}",
                            "objects": objects,
                            "notes": read_pptx_notes_xml(archive, slide_path),
                            "transition": xml_transition_payload(slide_root),
                            "timeline": xml_timeline_payload(slide_root),
                        },
                        index,
                        derive_summary=True,
                    )
                )

    return {
        "kind": "presentation",
        "format": "pptx",
        "path": path,
        **metadata,
        "warning": PRESENTATION_AUTOMATION_WARNING,
        "slides": slides,
    }


def pptx_placeholder_name(shape) -> str:
    if not getattr(shape, "is_placeholder", False):
        return ""
    try:
        placeholder_type = shape.placeholder_format.type
    except Exception:
        return ""
    name = getattr(placeholder_type, "name", None)
    return str(name or placeholder_type).upper()


def presentation_shape_paragraphs(shape) -> List[Dict[str, object]]:
    if not getattr(shape, "has_text_frame", False):
        return []
    paragraphs: List[Dict[str, object]] = []
    for paragraph in shape.text_frame.paragraphs:
        text = str(getattr(paragraph, "text", "") or "").strip()
        if not text:
            continue
        paragraphs.append(
            {
                "text": text,
                "level": max(int(getattr(paragraph, "level", 0) or 0), 0),
            }
        )
    return paragraphs


def presentation_table_cells(table) -> List[List[str]]:
    rows: List[List[str]] = []
    for row in getattr(table, "rows", []) or []:
        next_row: List[str] = []
        for cell in getattr(row, "cells", []) or []:
            next_row.append(str(getattr(cell, "text", "") or ""))
        if next_row:
            rows.append(next_row)
    return rows


def chart_series_payload(chart) -> List[Dict[str, object]]:
    payload: List[Dict[str, object]] = []
    try:
        series_iter = list(chart.series)
    except Exception:
        return payload
    for index, series in enumerate(series_iter, start=1):
        entry: Dict[str, object] = {"id": f"series:{index}"}
        name = str(getattr(series, "name", "") or "").strip()
        if name:
            entry["name"] = name
        try:
            values = list(getattr(series, "values"))
        except Exception:
            values = []
        if values:
            entry["values"] = values
        if len(entry) > 1:
            payload.append(entry)
    return payload


def chart_categories_payload(chart) -> List[str]:
    try:
        plots = list(chart.plots)
    except Exception:
        return []
    for plot in plots:
        try:
            categories = list(getattr(plot, "categories"))
        except Exception:
            categories = []
        if categories:
            return [str(value) for value in categories]
    return []


def presentation_shape_object(shape) -> Optional[Dict[str, object]]:
    object_id = f"shape:{getattr(shape, 'shape_id', '') or shape.name}"
    placeholder_kind = pptx_placeholder_name(shape)
    result: Dict[str, object] = {
        "id": object_id,
        "kind": "shape",
    }
    name = str(getattr(shape, "name", "") or "").strip()
    if name:
        result["name"] = name
    if placeholder_kind:
        result["placeholder_kind"] = placeholder_kind
    frame = emu_frame_dict(
        getattr(shape, "left", 0),
        getattr(shape, "top", 0),
        getattr(shape, "width", 0),
        getattr(shape, "height", 0),
    )
    result["frame"] = frame

    if getattr(shape, "has_text_frame", False):
        paragraphs = presentation_shape_paragraphs(shape)
        body = presentation_body_from_paragraphs(paragraphs)
        if placeholder_kind in PPTX_TITLE_PLACEHOLDER_TYPES:
            result["kind"] = "title"
            result["text"] = presentation_text_value(body)
        else:
            result["kind"] = "text_box"
            result["body"] = body
        return result

    if getattr(shape, "has_chart", False):
        result["kind"] = "chart"
        chart = getattr(shape, "chart", None)
        chart_type = getattr(chart, "chart_type", None)
        if chart_type is not None:
            result["chart_kind"] = str(getattr(chart_type, "name", chart_type))
        chart_title = ""
        if chart is not None:
            try:
                has_title = bool(getattr(chart, "has_title", False))
            except Exception:
                has_title = False
            if has_title:
                try:
                    chart_title = str(chart.chart_title.text_frame.text or "").strip()
                except Exception:
                    chart_title = ""
        if chart_title:
            result["title"] = chart_title
        categories = chart_categories_payload(chart) if chart is not None else []
        if categories:
            result["categories"] = categories
        series = chart_series_payload(chart) if chart is not None else []
        if series:
            result["series"] = series
        return result

    if getattr(shape, "has_table", False):
        result["kind"] = "table"
        table = getattr(shape, "table", None)
        if table is not None:
            result["table"] = {
                "rows": len(getattr(table, "rows", []) or []),
                "cols": len(getattr(table, "columns", []) or []),
            }
            cells = presentation_table_cells(table)
            if cells:
                result["table"]["cells"] = cells
        return result

    shape_type_name = str(getattr(getattr(shape, "shape_type", None), "name", "") or "").upper()
    if shape_type_name == "PICTURE":
        result["kind"] = "image"
        try:
            filename = str(getattr(shape.image, "filename", "") or "").strip()
            if filename:
                result["image_name"] = filename
        except Exception:
            pass
        alt_text = str(getattr(shape, "alt_text", "") or "").strip()
        if alt_text:
            result["alt_text"] = alt_text
    return result


def slide_layout_ref(slide) -> Optional[str]:
    layout = getattr(slide, "slide_layout", None)
    if layout is None:
        return None
    return normalize_layout_ref(getattr(layout, "name", None))


def slide_master_ref(slide) -> Optional[str]:
    layout = getattr(slide, "slide_layout", None)
    master = getattr(layout, "slide_master", None) if layout is not None else None
    return normalize_master_ref(getattr(master, "name", None))


def presentation_notes_from_slide(slide) -> List[str]:
    try:
        notes_slide = slide.notes_slide
    except Exception:
        return []
    notes: List[str] = []
    for shape in getattr(notes_slide, "shapes", []):
        if not getattr(shape, "has_text_frame", False):
            continue
        placeholder_kind = pptx_placeholder_name(shape)
        if placeholder_kind in PPTX_NOTES_IGNORED_PLACEHOLDER_TYPES:
            continue
        notes.extend(
            [str(item.get("text") or "") for item in presentation_shape_paragraphs(shape)]
        )
    return [note for note in notes if note.strip()]


def augment_presentation_document_from_xml(path: str, document: Dict[str, object]) -> Dict[str, object]:
    slides = document.get("slides")
    if not isinstance(slides, list) or not os.path.exists(path):
        return document
    sidecar_by_id: Dict[str, Dict[str, object]] = {}
    try:
        with zipfile.ZipFile(path, "r") as archive:
            presentation_xml = "ppt/presentation.xml"
            presentation_rels_xml = "ppt/_rels/presentation.xml.rels"
            if presentation_xml not in archive.namelist() or presentation_rels_xml not in archive.namelist():
                return document
            presentation_root = ET.fromstring(archive.read(presentation_xml))
            presentation_rels_root = ET.fromstring(archive.read(presentation_rels_xml))
            slide_targets: Dict[str, str] = {}
            for rel in presentation_rels_root:
                if local_name(rel.tag) != "Relationship":
                    continue
                rel_id = rel.attrib.get("Id")
                target = rel.attrib.get("Target")
                if rel_id and target:
                    slide_targets[rel_id] = normalize_relationship_target("ppt", target)
            slide_id_list = presentation_root.find(pptx_tag("sldIdLst"))
            if slide_id_list is None:
                return document
            for index, slide_ref in enumerate(slide_id_list.findall(pptx_tag("sldId")), start=1):
                slide_stable_id = str(slide_ref.attrib.get("id") or index)
                rel_id = slide_ref.attrib.get(f"{{{DOC_REL_NS}}}id")
                slide_path = slide_targets.get(rel_id or "")
                if not slide_path or slide_path not in archive.namelist():
                    continue
                slide_root = ET.fromstring(archive.read(slide_path))
                extra: Dict[str, object] = {}
                transition = xml_transition_payload(slide_root)
                if transition:
                    extra["transition"] = transition
                timeline = xml_timeline_payload(slide_root)
                if timeline:
                    extra["timeline"] = timeline
                if extra:
                    sidecar_by_id[f"slide:{slide_stable_id}"] = extra
    except Exception:
        return document

    if not sidecar_by_id:
        return document
    next_document = dict(document)
    next_slides = []
    for slide in slides:
        if not isinstance(slide, dict):
            next_slides.append(slide)
            continue
        merged = dict(slide)
        extra = sidecar_by_id.get(str(slide.get("id") or ""))
        if extra:
            merged.update(extra)
        next_slides.append(merged)
    next_document["slides"] = next_slides
    return next_document


def read_pptx_document_python(path: str, metadata: Dict[str, object]) -> Dict[str, object]:
    presentation = PptxPresentation(path)
    slides: List[Dict[str, object]] = []
    for index, slide in enumerate(presentation.slides, start=1):
        objects = []
        for shape in slide.shapes:
            obj = presentation_shape_object(shape)
            if obj is not None:
                objects.append(obj)
        slide_payload: Dict[str, object] = {
            "id": f"slide:{getattr(slide, 'slide_id', index)}",
            "objects": objects,
            "notes": presentation_notes_from_slide(slide),
        }
        layout_ref = slide_layout_ref(slide)
        if layout_ref:
            slide_payload["layout_ref"] = layout_ref
        master_ref = slide_master_ref(slide)
        if master_ref:
            slide_payload["master_ref"] = master_ref
        slides.append(normalize_presentation_slide_payload(slide_payload, index, derive_summary=True))
    return {
        "kind": "presentation",
        "format": "pptx",
        "path": path,
        **metadata,
        "warning": PRESENTATION_AUTOMATION_WARNING,
        "slides": slides,
    }


def read_pptx_document(path: str) -> Dict[str, object]:
    metadata = path_metadata(path)
    if not metadata["exists"]:
        return {
            "kind": "presentation",
            "format": "pptx",
            "path": path,
            **metadata,
            "warning": PRESENTATION_AUTOMATION_WARNING,
            "slides": [],
        }
    if python_pptx_available():
        try:
            return augment_presentation_document_from_xml(path, read_pptx_document_python(path, metadata))
        except Exception:
            pass
    return augment_presentation_document_from_xml(path, read_pptx_document_xml(path, metadata))


def choose_pptx_slide_layout(presentation, wants_body: bool, desired_layout_ref: Optional[str] = None):
    layouts = list(presentation.slide_layouts)
    if not layouts:
        raise RuntimeError("The presentation template does not provide any slide layouts.")
    preferred = None
    title_only = None
    desired = normalize_layout_ref(desired_layout_ref)
    for layout in layouts:
        layout_ref = normalize_layout_ref(getattr(layout, "name", None))
        placeholder_names = {pptx_placeholder_name(shape) for shape in layout.placeholders}
        has_title = bool(placeholder_names & {"TITLE", "CENTER_TITLE"})
        has_body = bool(placeholder_names & {"BODY", "OBJECT", "SUBTITLE"})
        if desired and layout_ref == desired:
            return layout
        if wants_body and has_title and has_body:
            return layout
        if has_title and title_only is None:
            title_only = layout
        if preferred is None:
            preferred = layout
    if wants_body:
        return title_only or preferred
    return title_only or preferred


def remove_pptx_slide(presentation, index: int) -> None:
    slide_id_list = presentation.slides._sldIdLst  # type: ignore[attr-defined]
    slide_id = slide_id_list[index]
    rel_id = slide_id.rId
    presentation.part.drop_rel(rel_id)
    del slide_id_list[index]


def clear_text_frame(text_frame) -> None:
    text_frame.clear()


def populate_text_frame(text_frame, paragraphs: List[object]) -> None:
    clear_text_frame(text_frame)
    normalized = body_items_from_outline(paragraphs)
    if not normalized:
        return
    first = text_frame.paragraphs[0]
    first.text = str(normalized[0].get("text") or "")
    first.level = int(normalized[0].get("level") or 0)
    for item in normalized[1:]:
        paragraph = text_frame.add_paragraph()
        paragraph.text = str(item.get("text") or "")
        paragraph.level = int(item.get("level") or 0)


def find_title_shape(slide):
    title_shape = getattr(slide.shapes, "title", None)
    if title_shape is not None and getattr(title_shape, "has_text_frame", False):
        return title_shape
    for shape in slide.shapes:
        if pptx_placeholder_name(shape) in {"TITLE", "CENTER_TITLE"} and getattr(shape, "has_text_frame", False):
            return shape
    return None


def find_body_shape(slide):
    for shape in slide.shapes:
        if pptx_placeholder_name(shape) in {"BODY", "OBJECT", "SUBTITLE"} and getattr(shape, "has_text_frame", False):
            return shape
    return None


def add_textbox_lines(slide, left, top, width, height, lines: List[str]):
    textbox = slide.shapes.add_textbox(left, top, width, height)
    populate_text_frame(
        textbox.text_frame,
        [{"text": line, "level": 0} for line in lines if str(line or "").strip()],
    )
    return textbox


def remove_shape(shape) -> None:
    element = getattr(shape, "_element", None)
    if element is None:
        return
    parent = element.getparent()
    if parent is None:
        return
    parent.remove(element)


def reorder_slide_shapes(slide, shapes: List[object]) -> None:
    sp_tree = getattr(slide.shapes, "_spTree", None)
    if sp_tree is None:
        return
    for shape in shapes:
        element = getattr(shape, "_element", None)
        if element is None:
            continue
        try:
            sp_tree.remove(element)
        except Exception:
            pass
        sp_tree.append(element)


def apply_frame_to_shape(shape, frame: object) -> None:
    normalized = normalize_presentation_frame(frame)
    if not normalized:
        return
    for source, attr in (("x", "left"), ("y", "top"), ("w", "width"), ("h", "height")):
        try:
            setattr(shape, attr, int(normalized.get(source) or 0))
        except Exception:
            continue


def shape_matches_payload(shape, payload: Dict[str, object], used_ids: set[str]):
    shape_id = f"shape:{getattr(shape, 'shape_id', '') or ''}"
    if shape_id in used_ids:
        return False
    placeholder_kind = pptx_placeholder_name(shape)
    desired_placeholder = str(payload.get("placeholder_kind") or "").upper()
    desired_kind = str(payload.get("kind") or "")
    if desired_placeholder and placeholder_kind == desired_placeholder:
        return True
    if desired_kind == "title" and placeholder_kind in PPTX_TITLE_PLACEHOLDER_TYPES:
        return True
    if desired_kind == "text_box" and placeholder_kind in PPTX_BODY_PLACEHOLDER_TYPES:
        return True
    return False


def safe_resolve_workspace_asset(raw_path: object) -> Optional[str]:
    text = str(raw_path or "").strip()
    if not text:
        return None
    try:
        resolved = resolve_workspace_path(text)
    except Exception:
        return None
    return resolved if os.path.exists(resolved) else None


def populate_table_shape(table, cells: List[List[str]]) -> None:
    for row_index, row in enumerate(cells):
        if row_index >= len(getattr(table, "rows", []) or []):
            break
        for col_index, value in enumerate(row):
            if col_index >= len(getattr(table, "columns", []) or []):
                break
            try:
                table.cell(row_index, col_index).text = str(value or "")
            except Exception:
                continue


def create_shape_from_payload(slide, payload: Dict[str, object]):
    kind = str(payload.get("kind") or "")
    frame = normalize_presentation_frame(payload.get("frame")) or emu_frame_dict(
        Inches(0.9),
        Inches(1.6),
        Inches(8.0),
        Inches(1.5),
    )
    if kind == "image":
        image_path = safe_resolve_workspace_asset(payload.get("image_ref"))
        if image_path:
            kwargs = {}
            if int(frame.get("w") or 0) > 0:
                kwargs["width"] = int(frame["w"])
            if int(frame.get("h") or 0) > 0:
                kwargs["height"] = int(frame["h"])
            return slide.shapes.add_picture(
                image_path,
                int(frame.get("x") or 0),
                int(frame.get("y") or 0),
                **kwargs,
            )
        return None
    if kind == "table":
        table_payload = payload.get("table") if isinstance(payload.get("table"), dict) else {}
        rows = max(int(table_payload.get("rows") or 0), 1)
        cols = max(int(table_payload.get("cols") or 0), 1)
        shape = slide.shapes.add_table(
            rows,
            cols,
            int(frame.get("x") or 0),
            int(frame.get("y") or 0),
            int(frame.get("w") or 0),
            int(frame.get("h") or 0),
        )
        cells = normalize_string_matrix(table_payload.get("cells"))
        if cells:
            populate_table_shape(shape.table, cells)
        return shape
    if kind in {"title", "text_box", "shape"} or "text" in payload or "body" in payload:
        textbox = slide.shapes.add_textbox(
            int(frame.get("x") or 0),
            int(frame.get("y") or 0),
            int(frame.get("w") or 0),
            int(frame.get("h") or 0),
        )
        text_body = payload.get("body")
        if "text" in payload and "body" not in payload:
            text_body = {"kind": "text", "text": str(payload.get("text") or "")}
        populate_text_frame(textbox.text_frame, presentation_paragraphs_from_body(text_body))
        return textbox
    return None


def update_notes_slide(slide, notes: object) -> None:
    normalized_notes = normalize_presentation_notes(notes)
    try:
        notes_slide = slide.notes_slide
    except Exception:
        return
    for shape in notes_slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        placeholder_kind = pptx_placeholder_name(shape)
        if placeholder_kind in PPTX_NOTES_IGNORED_PLACEHOLDER_TYPES:
            continue
        populate_text_frame(
            shape.text_frame,
            [{"text": line, "level": 0} for line in normalized_notes],
        )
        return


def replace_shape_with_payload(slide, shape, payload: Dict[str, object]):
    fallback_frame = emu_frame_dict(
        getattr(shape, "left", 0),
        getattr(shape, "top", 0),
        getattr(shape, "width", 0),
        getattr(shape, "height", 0),
    )
    next_payload = dict(payload)
    if "frame" not in next_payload:
        next_payload["frame"] = fallback_frame
    remove_shape(shape)
    return create_shape_from_payload(slide, next_payload)


def update_slide_from_payload(slide, slide_payload: Dict[str, object]) -> None:
    existing_shapes = {
        f"shape:{getattr(shape, 'shape_id', '') or shape.name}": shape for shape in slide.shapes
    }
    used_ids: set[str] = set()
    touched_title = False
    touched_body = False
    ordered_shapes: List[object] = []
    explicit_remove_ids = {
        str(value).strip()
        for value in (slide_payload.get("remove_object_ids") or [])
        if str(value or "").strip()
    }

    for obj in slide_payload.get("objects") or []:
        object_id = str(obj.get("id") or "")
        shape = existing_shapes.get(object_id)
        if shape is None:
            for candidate in slide.shapes:
                if shape_matches_payload(candidate, obj, used_ids):
                    shape = candidate
                    break
        if shape is None:
            shape = create_shape_from_payload(slide, obj)
        elif str(obj.get("kind") or "") == "image" and obj.get("image_ref"):
            shape = replace_shape_with_payload(slide, shape, obj)
        if shape is None:
            continue
        used_ids.add(f"shape:{getattr(shape, 'shape_id', '') or shape.name}")
        ordered_shapes.append(shape)
        apply_frame_to_shape(shape, obj.get("frame"))
        if getattr(shape, "has_table", False):
            table_payload = obj.get("table") if isinstance(obj.get("table"), dict) else {}
            cells = normalize_string_matrix(table_payload.get("cells"))
            if cells:
                populate_table_shape(shape.table, cells)
        if getattr(shape, "has_text_frame", False):
            body = obj.get("body")
            if "text" in obj and "body" not in obj:
                body = {"kind": "text", "text": str(obj.get("text") or "")}
            if body is not None:
                populate_text_frame(shape.text_frame, presentation_paragraphs_from_body(body))
        placeholder_kind = str(obj.get("placeholder_kind") or "").upper()
        kind = str(obj.get("kind") or "")
        if kind == "title" or placeholder_kind in PPTX_TITLE_PLACEHOLDER_TYPES:
            touched_title = True
        if placeholder_kind in PPTX_BODY_PLACEHOLDER_TYPES:
            touched_body = True

    for object_id in explicit_remove_ids:
        shape = existing_shapes.get(object_id)
        if shape is not None:
            remove_shape(shape)

    if slide_payload.get("prune_missing_objects") is True:
        for object_id, shape in existing_shapes.items():
            if object_id not in used_ids and object_id not in explicit_remove_ids:
                remove_shape(shape)

    if "title" in slide_payload and not touched_title:
        title_shape = find_title_shape(slide)
        if title_shape is not None:
            populate_text_frame(
                title_shape.text_frame,
                presentation_paragraphs_from_body({"kind": "text", "text": str(slide_payload.get("title") or "")}),
            )
        elif str(slide_payload.get("title") or "").strip():
            add_textbox_lines(
                slide,
                Inches(0.75),
                Inches(0.5),
                Inches(8.5),
                Inches(0.9),
                [str(slide_payload.get("title") or "")],
            )

    if "body" in slide_payload and not touched_body:
        body_shape = find_body_shape(slide)
        paragraphs = presentation_paragraphs_from_body(slide_payload.get("body"))
        if body_shape is not None:
            populate_text_frame(body_shape.text_frame, paragraphs)
        elif paragraphs:
            add_textbox_lines(
                slide,
                Inches(0.9),
                Inches(1.6),
                Inches(8.0),
                Inches(4.5),
                [str(item.get("text") or "") for item in paragraphs],
            )

    if "notes" in slide_payload:
        update_notes_slide(slide, slide_payload.get("notes"))

    if slide_payload.get("reorder_objects") is True and ordered_shapes:
        reorder_slide_shapes(slide, ordered_shapes)


def write_pptx_document(path: str, slides: List[Dict[str, object]], expected_etag: Optional[str]) -> Dict[str, object]:
    assert_expected_etag(path, expected_etag)
    normalized_slides = [
        normalize_presentation_slide_payload(slide, index, derive_summary=False)
        for index, slide in enumerate(slides, start=1)
    ]
    if not normalized_slides:
        raise RuntimeError("Presentation AIO must contain at least one slide.")
    if not python_pptx_available():
        raise RuntimeError(
            "Structured .pptx automation requires python-pptx. Rebuild the OpenClaw runtime image."
        )

    presentation = PptxPresentation(path) if os.path.exists(path) else PptxPresentation()
    existing_slides = {
        f"slide:{getattr(slide, 'slide_id', index)}": slide
        for index, slide in enumerate(presentation.slides, start=1)
    }

    for slide_payload in normalized_slides:
        slide_id = str(slide_payload.get("id") or "")
        ppt_slide = existing_slides.get(slide_id)
        if ppt_slide is None:
            body = slide_payload.get("body") or {"kind": "text", "text": ""}
            ppt_slide = presentation.slides.add_slide(
                choose_pptx_slide_layout(
                    presentation,
                    bool(presentation_paragraphs_from_body(body)),
                    slide_payload.get("layout_ref"),
                )
            )
        update_slide_from_payload(ppt_slide, slide_payload)

    temp_fd, temp_path = tempfile.mkstemp(prefix=".entropic-office-", dir=posixpath.dirname(path))
    os.close(temp_fd)
    try:
        presentation.save(temp_path)
        os.replace(temp_path, path)
    finally:
        try:
            if os.path.exists(temp_path):
                os.unlink(temp_path)
        except OSError:
            pass
    return read_pptx_document(path)


def spreadsheet_document_to_aio(document: Dict[str, object]) -> Dict[str, object]:
    workbook_id = f"workbook:{slugify_identifier(basename_without_extension(str(document.get('path') or 'workbook')), 'workbook')}"
    worksheets = []
    for index, sheet in enumerate(document.get("sheets") or [], start=1):
        if not isinstance(sheet, dict):
            continue
        sheet_name = str(sheet.get("name") or f"Sheet{index}").strip() or f"Sheet{index}"
        worksheet_slug = slugify_identifier(sheet_name, f"sheet{index}")
        cells = []
        for raw_cell in sheet.get("cells") or []:
            if not isinstance(raw_cell, dict):
                continue
            ref = str(raw_cell.get("ref") or "").strip().upper()
            row = int(raw_cell.get("row") or 0)
            col = int(raw_cell.get("col") or 0)
            if not ref:
                if row <= 0 or col <= 0:
                    continue
                ref = cell_ref(row, col)
            elif row <= 0 or col <= 0:
                row, col = parse_cell_ref(ref)
            formula = str(raw_cell.get("formula") or "").strip()
            value = str(raw_cell.get("value") or "")
            display = str(raw_cell.get("display") or value)
            kind = str(raw_cell.get("kind") or ("formula" if formula else "string"))
            entry = {
                "id": f"cell:{worksheet_slug}:{ref}",
                "kind": "cell",
                "ref": ref,
                "row": row,
                "col": col,
            }
            if formula:
                entry["formula"] = formula
                if value:
                    entry["value"] = value
                if display and display != value:
                    entry["display"] = display
            else:
                entry["value"] = value
                entry["value_kind"] = kind
                if display and display != value:
                    entry["display"] = display
            cells.append(entry)
        worksheets.append(
            {
                "id": f"worksheet:{worksheet_slug}",
                "kind": "worksheet",
                "name": sheet_name,
                "extent": {
                    "rows": int(sheet.get("rowCount") or 0),
                    "cols": int(sheet.get("colCount") or 0),
                },
                "cells": cells,
            }
        )
    return aio_envelope(
        "spreadsheet",
        document,
        {
            "id": workbook_id,
            "kind": "workbook",
            "worksheets": worksheets,
        },
        capability_set="current_entropic_office",
    )


def document_document_to_aio(document: Dict[str, object]) -> Dict[str, object]:
    document_id = f"document:{slugify_identifier(basename_without_extension(str(document.get('path') or 'document')), 'document')}"
    blocks = []
    for index, paragraph in enumerate(document.get("paragraphs") or [], start=1):
        blocks.append(
            {
                "id": f"block:{index}",
                "kind": "paragraph",
                "index": index,
                "text": str(paragraph),
            }
        )
    return aio_envelope(
        "document",
        document,
        {
            "id": document_id,
            "kind": "document",
            "blocks": blocks,
        },
        capability_set="current_entropic_office",
    )


def presentation_document_to_aio(document: Dict[str, object]) -> Dict[str, object]:
    deck_id = f"deck:{slugify_identifier(basename_without_extension(str(document.get('path') or 'deck')), 'deck')}"
    slides = []
    for index, slide in enumerate(document.get("slides") or [], start=1):
        normalized = normalize_presentation_slide_payload(slide, index, derive_summary=True)
        slide_entry: Dict[str, object] = {
            "id": str(normalized.get("id") or f"slide:{index}"),
            "kind": "slide",
            "index": index,
        }
        if normalized.get("layout_ref"):
            slide_entry["layout_ref"] = str(normalized.get("layout_ref"))
        if normalized.get("master_ref"):
            slide_entry["master_ref"] = str(normalized.get("master_ref"))
        title = str(normalized.get("title") or "").strip()
        if title:
            slide_entry["title"] = title
        body = normalized.get("body")
        if body is not None and presentation_paragraphs_from_body(body):
            slide_entry["body"] = body
        if normalized.get("objects"):
            slide_entry["objects"] = normalized.get("objects")
        if normalized.get("transition"):
            slide_entry["transition"] = normalized.get("transition")
        if normalized.get("timeline"):
            slide_entry["timeline"] = normalized.get("timeline")
        notes = normalize_presentation_notes(normalized.get("notes"))
        if notes:
            slide_entry["notes"] = notes
        slides.append(slide_entry)
    return aio_envelope(
        "presentation",
        document,
        {
            "id": deck_id,
            "kind": "deck",
            "slides": slides,
        },
        capability_set="current_entropic_office",
    )


def document_to_aio(document: Dict[str, object]) -> Dict[str, object]:
    kind = str(document.get("kind") or "").strip()
    if kind == "spreadsheet":
        return spreadsheet_document_to_aio(document)
    if kind == "document":
        return document_document_to_aio(document)
    if kind == "presentation":
        return presentation_document_to_aio(document)
    raise RuntimeError(f"Unsupported AIO conversion for `{kind or 'unknown'}`.")


def attach_aio_projection(document: Dict[str, object]) -> Dict[str, object]:
    result = dict(document)
    result["aio"] = document_to_aio(document)
    return result


def spreadsheet_payload_from_aio(payload: Dict[str, object]) -> Dict[str, object]:
    ensure_aio_payload(payload)
    if str(payload.get("kind") or "").strip() != "spreadsheet":
        raise RuntimeError("Expected a spreadsheet AIO object.")
    object_payload = payload.get("object")
    if not isinstance(object_payload, dict):
        raise RuntimeError("AIO spreadsheet object is missing.")
    worksheets = object_payload.get("worksheets")
    if not isinstance(worksheets, list):
        raise RuntimeError("AIO spreadsheet object must contain `worksheets`.")
    sheets = []
    for index, worksheet in enumerate(worksheets, start=1):
        if not isinstance(worksheet, dict):
            continue
        name = str(worksheet.get("name") or f"Sheet{index}").strip() or f"Sheet{index}"
        cells = []
        for raw_cell in worksheet.get("cells") or []:
            if not isinstance(raw_cell, dict):
                continue
            ref = str(raw_cell.get("ref") or "").strip().upper()
            row = int(raw_cell.get("row") or 0)
            col = int(raw_cell.get("col") or 0)
            if not ref:
                if row <= 0 or col <= 0:
                    continue
                ref = cell_ref(row, col)
            elif row <= 0 or col <= 0:
                row, col = parse_cell_ref(ref)
            formula = str(raw_cell.get("formula") or "").strip()
            value = str(raw_cell.get("value") or "")
            display = str(raw_cell.get("display") or value)
            value_kind = str(raw_cell.get("value_kind") or raw_cell.get("kind") or "").strip()
            cell_entry = {
                "ref": ref,
                "row": row,
                "col": col,
                "display": display,
            }
            if formula:
                cell_entry["formula"] = formula
                cell_entry["value"] = value
            else:
                cell_entry["value"] = value
                if value_kind:
                    cell_entry["kind"] = value_kind
            cells.append(cell_entry)
        sheets.append({"name": name, "cells": cells})
    return {
        "expectedEtag": extract_aio_source_etag(payload),
        "sheets": sheets,
    }


def document_payload_from_aio(payload: Dict[str, object]) -> Dict[str, object]:
    ensure_aio_payload(payload)
    if str(payload.get("kind") or "").strip() != "document":
        raise RuntimeError("Expected a document AIO object.")
    object_payload = payload.get("object")
    if not isinstance(object_payload, dict):
        raise RuntimeError("AIO document object is missing.")
    blocks = object_payload.get("blocks")
    if not isinstance(blocks, list):
        raise RuntimeError("AIO document object must contain `blocks`.")
    paragraphs: List[str] = []
    for block in blocks:
        if not isinstance(block, dict):
            continue
        block_kind = str(block.get("kind") or "paragraph").strip()
        if block_kind == "list":
            for item in block.get("items") or []:
                if isinstance(item, dict):
                    paragraphs.append(str(item.get("text") or ""))
                else:
                    paragraphs.append(str(item))
            continue
        paragraphs.append(str(block.get("text") or ""))
    return {
        "expectedEtag": extract_aio_source_etag(payload),
        "paragraphs": paragraphs,
    }


def presentation_payload_from_aio(payload: Dict[str, object]) -> Dict[str, object]:
    ensure_aio_payload(payload)
    if str(payload.get("kind") or "").strip() != "presentation":
        raise RuntimeError("Expected a presentation AIO object.")
    object_payload = payload.get("object")
    if not isinstance(object_payload, dict):
        raise RuntimeError("AIO presentation object is missing.")
    slides = object_payload.get("slides")
    if not isinstance(slides, list):
        raise RuntimeError("AIO presentation object must contain `slides`.")
    return {
        "expectedEtag": extract_aio_source_etag(payload),
        "slides": [
            normalize_presentation_slide_payload(slide, index, derive_summary=False)
            for index, slide in enumerate(slides, start=1)
        ],
    }


def inspect_aio(path: str) -> Dict[str, object]:
    extension = split_extension(path)
    if extension == ".csv":
        return spreadsheet_document_to_aio(read_csv_document(path))
    if extension == ".xlsx":
        return spreadsheet_document_to_aio(read_xlsx_document(path))
    if extension == ".docx":
        return document_document_to_aio(read_docx_document(path))
    if extension == ".pptx":
        return presentation_document_to_aio(read_pptx_document(path))
    if extension == ".xls":
        raise RuntimeError("Legacy .xls files are not supported yet. Save or convert the workbook as .xlsx.")
    raise RuntimeError("Unsupported format for AIO inspection.")


def apply_aio(path: str, payload: Dict[str, object]) -> Dict[str, object]:
    ensure_aio_payload(payload)
    payload_source_path = extract_aio_source_path(payload)
    if payload_source_path:
        try:
            normalized_payload_source_path = resolve_workspace_path(payload_source_path)
        except ValueError as error:
            raise RuntimeError(str(error))
        if posixpath.normpath(normalized_payload_source_path) != posixpath.normpath(path):
            raise RuntimeError("AIO payload source path does not match the requested path.")
    extension = split_extension(path)
    if extension in {".csv", ".xlsx"}:
        return save_spreadsheet(path, spreadsheet_payload_from_aio(payload))
    if extension == ".docx":
        return save_document(path, document_payload_from_aio(payload))
    if extension == ".pptx":
        return save_presentation(path, payload)
    if extension == ".xls":
        raise RuntimeError("Legacy .xls files are not supported yet. Save or convert the workbook as .xlsx.")
    raise RuntimeError("Unsupported format for AIO application.")


def inspect_spreadsheet(path: str) -> Dict[str, object]:
    extension = split_extension(path)
    if extension == ".csv":
        return attach_aio_projection(read_csv_document(path))
    if extension == ".xlsx":
        return attach_aio_projection(read_xlsx_document(path))
    if extension == ".xls":
        raise RuntimeError("Legacy .xls files are not supported yet. Save or convert the workbook as .xlsx.")
    raise RuntimeError("Unsupported spreadsheet format. Use .xlsx or .csv.")


def save_spreadsheet(path: str, payload: Dict[str, object]) -> Dict[str, object]:
    if isinstance(payload, dict) and payload.get("spec") == AIO_SPEC:
        payload = spreadsheet_payload_from_aio(payload)
    extension = split_extension(path)
    sheets = payload.get("sheets") or []
    expected_etag = payload.get("expectedEtag")
    if extension == ".csv":
        return attach_aio_projection(write_csv_document(path, sheets, expected_etag))
    if extension == ".xlsx":
        return attach_aio_projection(write_xlsx_document(path, sheets, expected_etag))
    if extension == ".xls":
        raise RuntimeError("Legacy .xls files are not supported yet. Save or convert the workbook as .xlsx.")
    raise RuntimeError("Unsupported spreadsheet format. Use .xlsx or .csv.")


def normalize_spreadsheet(path: str) -> Dict[str, object]:
    extension = split_extension(path)
    if extension == ".csv":
        return attach_aio_projection(read_csv_document(path))
    if extension == ".xlsx":
        document = read_xlsx_document(path)
        if document.get("exists") and legacy_minimal_xlsx(path):
            return attach_aio_projection(write_xlsx_document(path, document.get("sheets") or [], document.get("etag")))
        return attach_aio_projection(document)
    if extension == ".xls":
        raise RuntimeError("Legacy .xls files are not supported yet. Save or convert the workbook as .xlsx.")
    raise RuntimeError("Unsupported spreadsheet format. Use .xlsx or .csv.")


def inspect_document(path: str) -> Dict[str, object]:
    extension = split_extension(path)
    if extension == ".docx":
        return attach_aio_projection(read_docx_document(path))
    raise RuntimeError("Unsupported document format. Use .docx.")


def save_document(path: str, payload: Dict[str, object]) -> Dict[str, object]:
    if isinstance(payload, dict) and payload.get("spec") == AIO_SPEC:
        payload = document_payload_from_aio(payload)
    extension = split_extension(path)
    paragraphs = payload.get("paragraphs") or []
    expected_etag = payload.get("expectedEtag")
    if extension == ".docx":
        return attach_aio_projection(write_docx_document(path, [str(value) for value in paragraphs], expected_etag))
    raise RuntimeError("Unsupported document format. Use .docx.")


def save_presentation(path: str, payload: Dict[str, object]) -> Dict[str, object]:
    if isinstance(payload, dict) and payload.get("spec") == AIO_SPEC:
        payload = presentation_payload_from_aio(payload)
    extension = split_extension(path)
    slides = payload.get("slides") or []
    expected_etag = payload.get("expectedEtag")
    if extension == ".pptx":
        return attach_aio_projection(write_pptx_document(path, slides, expected_etag))
    raise RuntimeError("Unsupported presentation format. Use .pptx.")


def todo_spreadsheet_payload(items: Iterable[str]) -> Dict[str, object]:
    rows = [
        {"ref": "A1", "row": 1, "col": 1, "value": "Task", "kind": "string"},
        {"ref": "B1", "row": 1, "col": 2, "value": "Status", "kind": "string"},
    ]
    for index, item in enumerate(items, start=2):
        rows.append({"ref": cell_ref(index, 1), "row": index, "col": 1, "value": str(item), "kind": "string"})
        rows.append({"ref": cell_ref(index, 2), "row": index, "col": 2, "value": "Todo", "kind": "string"})
    return {"sheets": [{"name": "Sheet1", "cells": rows}]}


def blank_spreadsheet_payload() -> Dict[str, object]:
    return {"sheets": [{"name": "Sheet1", "cells": []}]}


def blank_document_payload(lines: Iterable[str]) -> Dict[str, object]:
    return {"paragraphs": [str(value) for value in lines]}


def cli_usage() -> str:
    return """Usage:
  # Preferred AIO workflow for .xlsx / .docx / .pptx
  entropic-office api inspect-aio <path>
  entropic-office api apply-aio <path>
  # Presentation AIO uses object.slides with sparse layout_ref/title/body summaries
  # plus objects[] entries carrying kind, frame, and optional notes/image_ref.

  # Legacy compatibility helpers
  entropic-office spreadsheet new <path>
  entropic-office spreadsheet todo <path> <item> [<item> ...]
  entropic-office document new <path>
  entropic-office document lines <path> <line> [<line> ...]

  # Low-level legacy APIs
  entropic-office api inspect-spreadsheet <path>
  entropic-office api normalize-spreadsheet <path>
  entropic-office api save-spreadsheet <path>
  entropic-office api inspect-document <path>
  entropic-office api save-document <path>
"""


def emit_json(payload: Dict[str, object]) -> None:
    sys.stdout.write(json.dumps(payload, ensure_ascii=False))
    sys.stdout.write("\n")


def run_api(argv: List[str]) -> Dict[str, object]:
    if len(argv) < 4:
        raise RuntimeError(cli_usage().strip())
    command = argv[2]
    path = resolve_workspace_path(argv[3])
    if command == "inspect-spreadsheet":
        return inspect_spreadsheet(path)
    if command == "normalize-spreadsheet":
        return normalize_spreadsheet(path)
    if command == "save-spreadsheet":
        payload = json.load(sys.stdin)
        return save_spreadsheet(path, payload)
    if command == "inspect-document":
        return inspect_document(path)
    if command == "save-document":
        payload = json.load(sys.stdin)
        return save_document(path, payload)
    if command == "inspect-aio":
        return inspect_aio(path)
    if command == "apply-aio":
        payload = json.load(sys.stdin)
        return apply_aio(path, payload)
    raise RuntimeError(cli_usage().strip())


def run_cli(argv: List[str]) -> Dict[str, object]:
    if len(argv) < 4:
        raise RuntimeError(cli_usage().strip())
    category = argv[1]
    command = argv[2]
    path = resolve_workspace_path(argv[3])
    if category == "spreadsheet" and command == "new":
        return save_spreadsheet(path, blank_spreadsheet_payload())
    if category == "spreadsheet" and command == "todo":
        items = argv[4:] or [f"Item {index}" for index in range(1, 11)]
        return save_spreadsheet(path, todo_spreadsheet_payload(items))
    if category == "document" and command == "new":
        return save_document(path, blank_document_payload([]))
    if category == "document" and command == "lines":
        return save_document(path, blank_document_payload(argv[4:]))
    raise RuntimeError(cli_usage().strip())


def main() -> int:
    try:
        if len(sys.argv) < 2:
            raise RuntimeError(cli_usage().strip())
        if sys.argv[1] == "api":
            result = run_api(sys.argv)
        else:
            result = run_cli(sys.argv)
        emit_json({"ok": True, "result": result})
        return 0
    except Exception as error:
        emit_json({"ok": False, "error": str(error)})
        return 1


if __name__ == "__main__":
    raise SystemExit(main())
